""""PPT file translation"""
import argparse

import os

from pptx import Presentation

from extension.files.utils import should_translate
from service.mt import translator_factory
from service.utils import detect_lang


def scan_doc(ppt, new_ppt):
    runs = []  # 待翻译对象
    for i, (slide, new_slide) in enumerate(zip(ppt.slides, new_ppt.slides)):  # 对每一页
        # print("Slide{}".format(i + 1), slide.slide_layout)
        for j, (shape, new_shape) in enumerate(zip(slide.shapes, new_slide.shapes)):  # 对每一对象
            # print("Shape{}".format(j + 1), shape)
            if shape.has_text_frame:  # 有文本内容
                text_frame = shape.text_frame
                for k, (paragraph, new_paragraph) in enumerate(zip(text_frame.paragraphs, new_shape.text_frame.paragraphs)):
                    # print("\tParagraph{}".format(k + 1), paragraph.text)
                    if should_translate(new_paragraph.text):
                        runs.append(new_paragraph)
            elif shape.has_table:  # 对表格
                table = shape.table
                new_table = new_shape.table
                for row, new_row in zip(table.rows, new_table.rows):
                    for cell, new_cell in zip(row.cells, new_row.cells):
                        if should_translate(new_cell.text_frame.text) :
                            runs.append(new_cell.text_frame)  # 翻译表格中文本
    return runs


def translate_ppt_auto(in_fn, source_lang="auto", target_lang="zh", translation_file=None, callbacker=None):
    paths = os.path.splitext(in_fn)
    docx_fn = in_fn

    if translation_file is None:
        translated_fn = paths[0] + "-translated.pptx"
    else:
        translated_fn = translation_file

    doc = Presentation(docx_fn)
    translated_doc = Presentation(docx_fn)  # 原样拷贝源文档
    runs = scan_doc(doc, translated_doc)  # 提取源文档中可以翻译的文本对象

    if source_lang == "auto":
        source_lang = detect_lang(runs[0].text)

    translator = translator_factory.get_translator(source_lang, target_lang)

    if translator is None:
        raise ValueError("给定语言不支持: {}".format(source_lang+"-"+target_lang))

    txt_list = [r.text for r in runs]
    result_list = translator.translate_list(txt_list, sl=source_lang, tl=target_lang, callbacker=callbacker)  # translate
    for i in range(len(runs)):
        runs[i].text = result_list[i]  # 替换源文本为翻译文本，其他不变（位置样式和非文本内容）

    translated_doc.save(translated_fn)

    return translated_fn


if __name__ == "__main__":
    arg_parser = argparse.ArgumentParser("PPT File Translator")
    arg_parser.add_argument("--to_lang", type=str, default="zh", help="target language")
    arg_parser.add_argument("--input_file", type=str, required=True, help="file to be translated")
    arg_parser.add_argument("--output_file", type=str, default=None, help="translation file")
    args = arg_parser.parse_args()

    to_lang = args.to_lang
    in_file = args.input_file
    out_file = args.output_file

    translated_fn = translate_ppt_auto(in_file, target_lang=to_lang, translation_file=out_file)

    import webbrowser

    webbrowser.open(in_file)
    webbrowser.open(translated_fn)

